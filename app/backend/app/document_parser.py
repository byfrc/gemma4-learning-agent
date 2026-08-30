from __future__ import annotations

import csv
import json
import os
import shutil
import subprocess
import tempfile
from dataclasses import dataclass
from pathlib import Path
from typing import Any

import numpy as np


SUPPORTED_DOCUMENT_SUFFIXES = {
    ".txt",
    ".md",
    ".csv",
    ".pdf",
    ".ppt",
    ".pptx",
}


class DocumentParseError(ValueError):
    """Raised when an uploaded document cannot be converted into text."""


@dataclass(frozen=True)
class ParsedSection:
    text: str
    location: str | None = None
    page: int | None = None
    slide: int | None = None
    document_type: str = "text"


def normalize_text(text: str) -> str:
    text = str(text or "")
    text = text.replace("\x00", " ")
    text = text.replace("\u200b", "").replace("\ufeff", "")
    return " ".join(text.split()).strip()


def _decode_bytes(raw: bytes) -> str:
    for encoding in ("utf-8-sig", "utf-8", "gb18030", "gbk"):
        try:
            return raw.decode(encoding)
        except UnicodeDecodeError:
            pass
    return raw.decode("utf-8", errors="ignore")


def _text_section(path: Path) -> list[ParsedSection]:
    text = normalize_text(_decode_bytes(path.read_bytes()))
    return [
        ParsedSection(
            text=text,
            document_type=path.suffix.lower().lstrip(".") or "text",
        )
    ] if text else []


def _csv_section(path: Path) -> list[ParsedSection]:
    rows = []
    decoded = _decode_bytes(path.read_bytes())
    for row in csv.reader(decoded.splitlines()):
        text = "；".join(
            f"字段{i + 1}：{cell.strip()}"
            for i, cell in enumerate(row)
            if cell.strip()
        )
        if text:
            rows.append(text)

    text = normalize_text("\n".join(rows))
    return [
        ParsedSection(text=text, document_type="csv"),
    ] if text else []


class _PaddleOCREngine:
    def __init__(self, lang: str, device: str):
        try:
            from paddleocr import PaddleOCR
        except ImportError as exc:
            raise DocumentParseError(
                "该 PDF 可能是扫描件，当前未安装 PaddleOCR。"
                "请安装 OCR 依赖，或将 DOCUMENT_OCR_ENABLED 设置为 false。"
            ) from exc

        self.engine = None
        modern_kwargs = {
            "lang": lang,
            "use_doc_orientation_classify": False,
            "use_doc_unwarping": False,
            "use_textline_orientation": False,
        }
        if device:
            modern_kwargs["device"] = device

        try:
            self.engine = PaddleOCR(**modern_kwargs)
        except TypeError:
            legacy_kwargs = {"lang": lang}
            if device:
                legacy_kwargs["use_gpu"] = device.startswith("gpu")
            self.engine = PaddleOCR(**legacy_kwargs)

    @staticmethod
    def _result_payload(result: Any) -> dict[str, Any]:
        payload = result
        json_value = getattr(result, "json", None)
        if json_value is not None:
            payload = json_value() if callable(json_value) else json_value
        if isinstance(payload, str):
            try:
                payload = json.loads(payload)
            except json.JSONDecodeError:
                return {}
        if not isinstance(payload, dict):
            return {}
        nested = payload.get("res")
        return nested if isinstance(nested, dict) else payload

    def extract(self, image: np.ndarray) -> str:
        try:
            results = self.engine.predict(image)
            texts: list[str] = []
            for result in results or []:
                payload = self._result_payload(result)
                values = payload.get("rec_texts", [])
                if isinstance(values, list):
                    texts.extend(str(value) for value in values if str(value).strip())
            text = normalize_text(" ".join(texts))
            if text:
                return text
        except AttributeError:
            pass

        # Compatibility fallback for older PaddleOCR releases.
        try:
            result = self.engine.ocr(image, cls=True)
        except Exception as exc:
            raise DocumentParseError(f"OCR 处理页面失败：{exc}") from exc

        texts = []
        for line_group in result or []:
            for line in line_group or []:
                if len(line) >= 2 and isinstance(line[1], (list, tuple)):
                    texts.append(str(line[1][0]))
        return normalize_text(" ".join(texts))


class DocumentParser:
    def __init__(
        self,
        ocr_mode: str = "auto",
        ocr_lang: str = "ch",
        ocr_device: str = "cpu",
        office_converter: str = "soffice",
        office_timeout_seconds: int = 120,
    ):
        self.ocr_mode = (ocr_mode or "auto").strip().lower()
        self.ocr_lang = ocr_lang
        self.ocr_device = ocr_device
        self.office_converter = office_converter
        self.office_timeout_seconds = office_timeout_seconds
        self._ocr_engine: _PaddleOCREngine | None = None
        self._ocr_unavailable = False

    def parse(self, path: Path) -> list[ParsedSection]:
        suffix = path.suffix.lower()
        if suffix not in SUPPORTED_DOCUMENT_SUFFIXES:
            raise DocumentParseError(
                f"暂不支持 {suffix or '该文件'} 格式。"
            )
        if not path.exists():
            raise DocumentParseError("文档文件不存在。")

        if suffix in {".txt", ".md"}:
            return _text_section(path)
        if suffix == ".csv":
            return _csv_section(path)
        if suffix == ".pdf":
            return self._parse_pdf(path)
        if suffix == ".pptx":
            return self._parse_pptx(path, source_type="pptx")
        return self._parse_legacy_ppt(path)

    def _parse_pdf(self, path: Path) -> list[ParsedSection]:
        try:
            import fitz
        except ImportError as exc:
            raise DocumentParseError(
                "PDF 解析需要 PyMuPDF，请先安装后端依赖。"
            ) from exc

        sections = []
        try:
            document = fitz.open(str(path))
        except Exception as exc:
            raise DocumentParseError(f"无法打开 PDF：{exc}") from exc

        try:
            for page_number, page in enumerate(document, start=1):
                text = normalize_text(page.get_text("text", sort=True))
                if not text and self._ocr_allowed:
                    text = self._ocr_page(page)
                if text:
                    sections.append(
                        ParsedSection(
                            text=text,
                            location=f"第 {page_number} 页",
                            page=page_number,
                            document_type="pdf",
                        )
                    )
        finally:
            document.close()

        return sections

    @property
    def _ocr_allowed(self) -> bool:
        return self.ocr_mode in {"1", "true", "yes", "on", "auto", "required"}

    @property
    def _ocr_required(self) -> bool:
        return self.ocr_mode in {"1", "true", "yes", "on", "required"}

    def _get_ocr_engine(self) -> _PaddleOCREngine | None:
        if self._ocr_engine is not None or self._ocr_unavailable:
            return self._ocr_engine
        try:
            self._ocr_engine = _PaddleOCREngine(
                lang=self.ocr_lang,
                device=self.ocr_device,
            )
        except DocumentParseError:
            if self._ocr_required:
                raise
            self._ocr_unavailable = True
            return None
        return self._ocr_engine

    def _ocr_page(self, page: Any) -> str:
        engine = self._get_ocr_engine()
        if engine is None:
            return ""

        try:
            import fitz

            pixmap = page.get_pixmap(
                matrix=fitz.Matrix(2, 2),
                alpha=False,
            )
            image = np.frombuffer(
                pixmap.samples,
                dtype=np.uint8,
            ).reshape(pixmap.height, pixmap.width, pixmap.n)
            if pixmap.n == 4:
                image = image[:, :, :3]
            return engine.extract(image)
        except DocumentParseError:
            raise
        except Exception as exc:
            raise DocumentParseError(f"OCR 处理 PDF 页面失败：{exc}") from exc

    @staticmethod
    def _shape_text_lines(shape: Any) -> list[str]:
        lines = []
        nested_shapes = getattr(shape, "shapes", None)
        if nested_shapes is not None:
            for child in nested_shapes:
                lines.extend(DocumentParser._shape_text_lines(child))

        if getattr(shape, "has_table", False):
            table_rows = []
            for row in shape.table.rows:
                cells = [
                    normalize_text(cell.text)
                    for cell in row.cells
                ]
                if any(cells):
                    table_rows.append(" | ".join(cells))
            if table_rows:
                lines.append("表格：" + "；".join(table_rows))
            return lines

        if getattr(shape, "has_text_frame", False):
            text = normalize_text(shape.text)
            if text:
                lines.append(text)
        return lines

    def _parse_pptx(
        self,
        path: Path,
        source_type: str,
    ) -> list[ParsedSection]:
        try:
            from pptx import Presentation
        except ImportError as exc:
            raise DocumentParseError(
                "PPT/PPTX 解析需要 python-pptx，请先安装后端依赖。"
            ) from exc

        try:
            presentation = Presentation(str(path))
        except Exception as exc:
            raise DocumentParseError(f"无法打开 PPTX：{exc}") from exc

        sections = []
        for slide_number, slide in enumerate(presentation.slides, start=1):
            lines = []
            for shape in slide.shapes:
                lines.extend(self._shape_text_lines(shape))

            try:
                notes_slide = slide.notes_slide
            except Exception:
                notes_slide = None
            if notes_slide is not None:
                notes_lines = []
                for shape in notes_slide.shapes:
                    notes_lines.extend(self._shape_text_lines(shape))
                if notes_lines:
                    lines.append("讲者备注：" + "；".join(notes_lines))

            deduplicated = []
            seen = set()
            for line in lines:
                normalized = normalize_text(line)
                if normalized and normalized not in seen:
                    deduplicated.append(normalized)
                    seen.add(normalized)

            text = normalize_text("\n".join(deduplicated))
            if text:
                sections.append(
                    ParsedSection(
                        text=text,
                        location=f"第 {slide_number} 张幻灯片",
                        slide=slide_number,
                        document_type=source_type,
                    )
                )
        return sections

    def _parse_legacy_ppt(self, path: Path) -> list[ParsedSection]:
        converter = (
            os.getenv("OFFICE_CONVERTER", "").strip()
            or self.office_converter
        )
        converter_path = shutil.which(converter) or (
            converter if Path(converter).exists() else None
        )

        if converter_path:
            with tempfile.TemporaryDirectory(prefix="gemma4-ppt-") as directory:
                try:
                    result = subprocess.run(
                        [
                            converter_path,
                            "--headless",
                            "--convert-to",
                            "pptx",
                            "--outdir",
                            directory,
                            str(path),
                        ],
                        capture_output=True,
                        text=True,
                        timeout=self.office_timeout_seconds,
                        check=False,
                    )
                except subprocess.TimeoutExpired as exc:
                    raise DocumentParseError(
                        f"旧版 PPT 转换超时（超过 {self.office_timeout_seconds} 秒）。"
                    ) from exc
                if result.returncode != 0:
                    detail = normalize_text(result.stderr or result.stdout)
                    raise DocumentParseError(
                        f"旧版 PPT 转换失败：{detail or 'LibreOffice 未返回具体原因。'}"
                    )

                converted = Path(directory) / f"{path.stem}.pptx"
                if not converted.exists():
                    candidates = list(Path(directory).glob("*.pptx"))
                    converted = candidates[0] if candidates else converted
                if converted.exists():
                    return self._parse_pptx(converted, source_type="ppt")

        catppt = shutil.which("catppt")
        if catppt:
            try:
                result = subprocess.run(
                    [catppt, str(path)],
                    capture_output=True,
                    timeout=self.office_timeout_seconds,
                    check=False,
                )
            except subprocess.TimeoutExpired as exc:
                raise DocumentParseError("旧版 PPT 文本提取超时。") from exc
            if result.returncode == 0:
                text = normalize_text(_decode_bytes(result.stdout))
                if text:
                    return [
                        ParsedSection(
                            text=text,
                            location="旧版 PPT 文本",
                            document_type="ppt",
                        )
                    ]

        raise DocumentParseError(
            "旧版 .ppt 需要 LibreOffice（soffice/libreoffice）或 catppt。"
            "建议在服务器安装 LibreOffice 后重新上传。"
        )
